import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation(output_path):
    prs = Presentation()

    # Helper function to add a slide with a title and content
    def add_slide(title_text, content_text=None, points=None):
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        if content_text:
            tf = slide.placeholders[1].text_frame
            tf.text = content_text
        
        if points:
            tf = slide.placeholders[1].text_frame
            tf.clear()  # Clear existing content if any
            for point in points:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
                if isinstance(point, list): # Handle sub-points if needed (not used here yet)
                    pass

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Multimodal Graph RAG"
    subtitle.text = "The Future of Contextual & Knowledge-Aware Retrieval\nPresented by: Project Team"

    # 2. Introduction: Why Multimodal Graph RAG?
    add_slide(
        "Why Multimodal Graph RAG?",
        points=[
            "Standard RAG is limited to flat vector retrieval.",
            "Graph RAG adds semantic relationships and domain context.",
            "Multimodal support enables searching images, PDFs, and text seamlessly.",
            "Combines Knowledge Graphs (Neo4j) with Vector DBs (Weaviate)."
        ]
    )

    # 3. Project Architecture
    add_slide(
        "Full-Stack Architecture",
        points=[
            "Frontend: Next.js (Modern, Responsive UI)",
            "Backend: FastAPI (High-performance Python API)",
            "Graph Database: Neo4j (Relationship Management)",
            "Vector Database: Weaviate (Similarity Search)",
            "LLM Engine: Google Gemini (Multimodal reasoning)"
        ]
    )

    # 4. Modalities Supported
    add_slide(
        "Multimodal Capabilities",
        points=[
            "Text: Natural language queries and semantic matching.",
            "Images: Direct image search and visual similarity.",
            "Documents/PDFs: Automated parsing and indexing of structured data.",
            "Cross-modal: Search text via images and vice-versa."
        ]
    )

    # 5. The Power of Graph RAG (Neo4j)
    add_slide(
        "Graph-Enhanced Retrieval",
        points=[
            "Nodes represent Entities (Products, Categories, Users).",
            "Edges represent Relationships (BELONGS_TO, RECOMMENDED_WITH).",
            "Retrieval uses Cypher queries to fetch deep contextual neighbors.",
            "Eliminates the 'hallucination' risk by grounding LLMs in facts."
        ]
    )

    # 6. Vector Search (Weaviate)
    add_slide(
        "Vector-Based Discovery",
        points=[
            "High-dimensional embeddings for semantic similarity.",
            "Supports fast retrieval for unstructured data.",
            "Weaviate integration allows scaling to millions of objects.",
            "Fallback mechanisms ensure reliability when services are unstable."
        ]
    )

    # 7. Gemini Integration
    add_slide(
        "AI Reasoning with Gemini",
        points=[
            "Gemini Pro Vision for image understanding.",
            "Gemini Pro for conversational reasoning and synthesis.",
            "Zero-shot and few-shot capabilities for diverse tasks.",
            "Unified API for all multimodal interactions."
        ]
    )

    # 8. User Interface & Experience
    add_slide(
        "Premium UI/UX",
        points=[
            "Built with Next.js and Tailwind CSS.",
            "Interactive search dashboard.",
            "Real-time chat interface with context history.",
            "Seamless file and image upload workflows."
        ]
    )

    # 9. Key Use Cases
    add_slide(
        "Industrial Use Cases",
        points=[
            "E-commerce: Smart product recommendations based on usage patterns.",
            "Enterprise Search: Navigating complex document hierarchies.",
            "Visual Search: Finding products or documents via screenshots.",
            "Knowledge Management: Linking siloed data points through graphs."
        ]
    )

    # 10. Conclusion & Future Scope
    add_slide(
        "Conclusion",
        points=[
            "Successfully combined Graph and Vector technologies.",
            "Robust multimodal processing pipeline.",
            "Future: Real-time graph updates and more modalities (Video/Audio).",
            "Questions? Thank you!"
        ]
    )

    # Save the presentation
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    output_file = "Multimodal_Graph_RAG_Presentation.pptx"
    create_presentation(output_file)
